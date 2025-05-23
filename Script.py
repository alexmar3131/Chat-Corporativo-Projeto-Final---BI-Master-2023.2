#Bloco1

# Script do Chat Eletronuclear — Bloco 1
# Funcionalidades: Upload, Extração, Banco Access, Relatórios e Interface

import tkinter as tk
from tkinter import filedialog
import os
import shutil
import pyodbc
import pandas as pd
import re
from datetime import datetime
from PIL import Image, ImageTk
from docx import Document
from docx.shared import Inches
import matplotlib.pyplot as plt
from PyPDF2 import PdfReader
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


# Caminhos fixos do sistema
CAMINHO_BANCO = r"D:\visita\Documents\Sistema TCC\Chat Eletronuclear\ReembolsosEscolares.accdb"
PASTA_RELATORIOS = r"D:\visita\Documents\Sistema TCC\Chat Eletronuclear\Relatórios"
PASTA_BACKUP = r"D:\visita\Documents\Sistema TCC\Chat Eletronuclear\BackUp Excel"
POPPLER_PATH = r"C:\Program Files\poppler-24.02.0\Library\bin"
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
LOGO_PATH = r"D:\visita\Documents\Sistema TCC\Chat Eletronuclear\Logo\Logo_Eletronuclear.png"

# Variáveis de controle
fase = "esperando_matricula"
nome_usuario = ""
matricula_usuario = ""
permissao_usuario = ""
status_usuario = ""
nome_funcionario_completo = ""
arquivos_reembolso = {}

# Bloco 2

def extrair_texto_pdf(caminho):
    try:
        reader = PdfReader(caminho)
        texto = '\n'.join(p.extract_text() or '' for p in reader.pages)
        if texto.strip():
            return texto.lower()
    except:
        pass
    try:
        imagens = convert_from_path(caminho, poppler_path=POPPLER_PATH)
        return '\n'.join(pytesseract.image_to_string(img, lang='por') for img in imagens).lower()
    except Exception as e:
        print('[ERRO OCR]', e)
        return ''

def classificar_por_conteudo(caminho_pdf):
    texto = extrair_texto_pdf(caminho_pdf)
    if "valor pago" in texto:
        return "comprovante"
    elif any(p in texto for p in ["nosso número", "ficha de compensação", "vencimento", "valor do documento"]):
        return "boleto"
    elif any(p in texto for p in ["nome do dependente", "ra", "etapa", "turno", "reembolso escolar", "nome escola"]):
        return "formulario"
    return "desconhecido"

def autenticar_matricula(matricula):
    global nome_usuario, matricula_usuario, permissao_usuario, status_usuario, fase, nome_funcionario_completo
    try:
        conn = pyodbc.connect(f"DRIVER={{Microsoft Access Driver (*.mdb, *.accdb)}};DBQ={CAMINHO_BANCO};")
        cursor = conn.cursor()
        cursor.execute("SELECT Nome_Funcionario, Permissao, Status FROM Funcionarios WHERE Matricula = ?", (matricula,))
        resultado = cursor.fetchone()
        conn.close()
        if resultado:
            nome_funcionario_completo, permissao_usuario, status_usuario = resultado
            nome_usuario = nome_funcionario_completo.split()[0]
            matricula_usuario = matricula
            texto_chat.config(state="normal")
            texto_chat.insert(tk.END, f"[Bot] Olá {nome_usuario}, perfil: {permissao_usuario} — Status: {status_usuario}\n")
            texto_chat.config(state="disabled")
            if status_usuario.lower() == "ativo":
                mostrar_menu()
            else:
                texto_chat.config(state="normal")
                texto_chat.insert(tk.END, "[Bot] Sua matrícula está inativa.\n")
                texto_chat.config(state="disabled")
        else:
            texto_chat.config(state="normal")
            texto_chat.insert(tk.END, "[Bot] Matrícula inexistente.\n")
            texto_chat.config(state="disabled")
    except Exception as e:
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, f"[Bot] Erro: {e}\n")
        texto_chat.config(state="disabled")

def mostrar_menu():
    texto_chat.config(state="normal")
    texto_chat.insert(tk.END, f"[Bot] {nome_usuario}, selecione uma opção:\n")
    texto_chat.insert(tk.END, "- [1] Solicitar reembolso\n")
    texto_chat.insert(tk.END, "- [2] Consultar reembolsos\n")
    if permissao_usuario.lower() == "administrador":
        texto_chat.insert(tk.END, "- [3] Gerar relatório Word\n")
        texto_chat.insert(tk.END, "- [4] Gerar relatório PDF\n")
        texto_chat.insert(tk.END, "- [5] Gerar apresentação PowerPoint\n")
        
    texto_chat.config(state="disabled")
    texto_chat.see(tk.END)

# Bloco 3

def extrair_dados():
    dados = {}

    caminho_comprovante = arquivos_reembolso.get("comprovante", "")
    texto_comp = extrair_texto_pdf(caminho_comprovante)
    dados["caminho_comprovante"] = caminho_comprovante
    match_valor = re.search(r"valor\s*pago[:\s]*r?\$?\s*([\d\.,]+)", texto_comp, re.IGNORECASE)
    if not match_valor:
        match_valor = re.search(r"(total|valor)\s*(pago|a pagar)[:\s]*r?\$?\s*([\d\.,]+)", texto_comp, re.IGNORECASE)
    valor_raw = match_valor.group(1 if match_valor.lastindex == 1 else 3) if match_valor else None
    dados["valor_pago"] = valor_raw.replace(".", "").replace(",", ".") if valor_raw else None

    caminho_boleto = arquivos_reembolso.get("boleto", "")
    texto_bol = extrair_texto_pdf(caminho_boleto)
    match_data = re.search(r"(vencimento)?[:\s]*?(\d{2}/\d{2}/\d{4})", texto_bol)
    if match_data:
        partes = match_data.group(2).split("/")
        dados["data_registro"] = match_data.group(2)
        dados["mes_competencia"] = int(partes[1])
        dados["ano_competencia"] = int(partes[2])
    else:
        dados["data_registro"] = ""
        dados["mes_competencia"] = 0
        dados["ano_competencia"] = 0

    caminho_form = arquivos_reembolso.get("formulario", "")
    texto_form = extrair_texto_pdf(caminho_form)
    linhas = texto_form.splitlines()
    for i, linha in enumerate(linhas):
        l = linha.lower().strip()
        if "nome  de dependente" in l and i + 1 < len(linhas):
            dados["nome_dependente"] = linhas[i + 1].strip()
        elif "código  de cadastramento" in l and i + 1 < len(linhas):
            valores = linhas[i + 1].strip().split()
            if len(valores) >= 3:
                try:
                    dados["codigo_dependente"] = int(valores[0])
                except:
                    dados["codigo_dependente"] = 0
                dados["serie_escolar"] = " ".join(valores[1:-1])
        elif "documento  apresentado" in l and i + 1 < len(linhas):
            linha_doc = linhas[i + 1].lower()
            if "mensalidade" in linha_doc:
                dados["tipo_documento"] = "Mensalidade"
            elif "matrícula" in linha_doc:
                dados["tipo_documento"] = "Matrícula"
            elif "material" in linha_doc:
                dados["tipo_documento"] = "Material"
            else:
                dados["tipo_documento"] = "Outro"

    dados.setdefault("nome_dependente", "")
    dados.setdefault("codigo_dependente", 0)
    dados.setdefault("serie_escolar", "")
    dados.setdefault("tipo_documento", "")
    dados["nome_funcionario"] = nome_funcionario_completo

    print("\n[DADOS EXTRAÍDOS]\n", dados)
    return dados

def inserir_dados_no_access(dados):
    try:
        if not matricula_usuario:
            print("[ERRO INSERT] Matrícula do usuário não está definida.")
            return

        conn = pyodbc.connect(f"DRIVER={{Microsoft Access Driver (*.mdb, *.accdb)}};DBQ={CAMINHO_BANCO};")
        cursor = conn.cursor()

        cursor.execute("""
            SELECT COUNT(*) FROM ReembolsosEscolares
            WHERE Matricula = ? AND Nome_Dependente = ? AND Mes_Competencia = ? AND Ano_Competencia = ?
        """, (
            matricula_usuario,
            dados.get("nome_dependente", ""),
            dados.get("mes_competencia", 0),
            dados.get("ano_competencia", 0)
        ))
        if cursor.fetchone()[0] > 0:
            print("[AVISO] Registro já existe. Inserção ignorada.")
            conn.close()
            return

        insert = """
            INSERT INTO ReembolsosEscolares (
                Matricula, Nome_Funcionario, Nome_Dependente, Codigo_Dependente, Serie_Escolar,
                Tipo_Documento, Mes_Competencia, Ano_Competencia, Valor_Pago,
                Caminho_Comprovante, Data_Registro
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """
        valores = (
            matricula_usuario,
            dados.get("nome_funcionario", ""),
            dados.get("nome_dependente", ""),
            dados.get("codigo_dependente", 0),
            dados.get("serie_escolar", ""),
            dados.get("tipo_documento", ""),
            dados.get("mes_competencia", 0),
            dados.get("ano_competencia", 0),
            float(dados.get("valor_pago", 0)),
            dados.get("caminho_comprovante", ""),
            dados.get("data_registro", "")
        )
        cursor.execute(insert, valores)
        conn.commit()
        conn.close()
        print("[INSERT OK] Registro inserido com sucesso.")

    except Exception as e:
        print("[ERRO INSERT]", e)

def gerar_backup_excel():
    try:
        conn = pyodbc.connect(f"DRIVER={{Microsoft Access Driver (*.mdb, *.accdb)}};DBQ={CAMINHO_BANCO};")
        df1 = pd.read_sql("SELECT * FROM ReembolsosEscolares", conn)
        df2 = pd.read_sql("SELECT * FROM Funcionarios", conn)
        conn.close()

        if not os.path.exists(PASTA_BACKUP):
            os.makedirs(PASTA_BACKUP)

        nome_arquivo = os.path.join(PASTA_BACKUP, f"Backup_{datetime.now():%Y%m%d_%H%M%S}.xlsx")
        with pd.ExcelWriter(nome_arquivo) as writer:
            df1.to_excel(writer, sheet_name="ReembolsosEscolares", index=False)
            df2.to_excel(writer, sheet_name="Funcionarios", index=False)
        print(f"[BACKUP EXCEL GERADO] {nome_arquivo}")

    except Exception as e:
        print("[ERRO BACKUP]", e)

# Bloco 4

def gerar_grafico_valores(df):
    df = df.copy()
    df["Valor_Pago"] = df["Valor_Pago"].astype(str).str.replace(",", ".").astype(float)
    df_val = df.groupby("Ano_Competencia")["Valor_Pago"].sum().reset_index()
    plt.figure(figsize=(6, 4))
    plt.bar(df_val["Ano_Competencia"].astype(str), df_val["Valor_Pago"], color="orange")
    plt.title("Valor Reembolsado por Ano")
    plt.xlabel("Ano")
    plt.ylabel("Valor Total (R$)")
    plt.tight_layout()
    caminho = os.path.join(PASTA_RELATORIOS, "grafico_valor_reembolsado_por_ano.png")
    plt.savefig(caminho)
    plt.close()
    return caminho

def gerar_grafico_funcionarios(df):
    df_func = df.groupby("Ano_Competencia")["Matricula"].nunique().reset_index()
    plt.figure(figsize=(6, 4))
    plt.bar(df_func["Ano_Competencia"].astype(str), df_func["Matricula"], color="orange")
    plt.title("Funcionários Reembolsados por Ano")
    plt.xlabel("Ano")
    plt.ylabel("Quantidade")
    plt.tight_layout()
    caminho = os.path.join(PASTA_RELATORIOS, "grafico_funcionarios_reembolsados_por_ano.png")
    plt.savefig(caminho)
    plt.close()
    return caminho

def gerar_relatorio_word_gerencial():
    try:
        conn = pyodbc.connect(f"DRIVER={{Microsoft Access Driver (*.mdb, *.accdb)}};DBQ={CAMINHO_BANCO};")
        df_reemb = pd.read_sql("SELECT * FROM ReembolsosEscolares", conn)
        df_func = pd.read_sql("SELECT * FROM Funcionarios", conn)
        conn.close()

        df_reemb["Ano_Competencia"] = pd.to_numeric(df_reemb["Ano_Competencia"], errors="coerce").astype("Int64")
        df_reemb["Valor_Pago"] = df_reemb["Valor_Pago"].astype(str).str.replace(",", ".").astype(float)
        ativos_df = df_func[df_func["Status"].str.lower() == "ativo"]

        doc = Document()
        doc.add_heading("Relatório Gerencial de Reembolsos", 0)
        doc.add_heading("Resumo de Funcionários", level=1)
        doc.add_paragraph(f"Total: {len(df_func)}")
        doc.add_paragraph(f"Ativos: {len(ativos_df)}")
        doc.add_paragraph(f"Inativos: {len(df_func) - len(ativos_df)}")

        doc.add_heading("Totais por Ano (Todos)", level=2)
        agrupado_todos = df_reemb.groupby("Ano_Competencia").agg({"Matricula": "nunique", "Valor_Pago": "sum"}).reset_index()
        table = doc.add_table(rows=1, cols=3)
        for i, text in enumerate(["Ano", "Funcionários", "Valor Total"]):
            table.rows[0].cells[i].text = text
        for _, row in agrupado_todos.iterrows():
            linha = table.add_row().cells
            linha[0].text = str(int(row["Ano_Competencia"]))
            linha[1].text = str(int(row["Matricula"]))
            linha[2].text = f"R$ {float(row['Valor_Pago']):,.2f}"

        doc.add_heading("Totais por Ano (Ativos)", level=2)
        df_ativos = df_reemb[df_reemb["Matricula"].isin(ativos_df["Matricula"])]
        agrupado_ativos = df_ativos.groupby("Ano_Competencia").agg({"Matricula": "nunique", "Valor_Pago": "sum"}).reset_index()
        table2 = doc.add_table(rows=1, cols=3)
        for i, text in enumerate(["Ano", "Funcionários Ativos", "Valor Ativos"]):
            table2.rows[0].cells[i].text = text
        for _, row in agrupado_ativos.iterrows():
            linha = table2.add_row().cells
            linha[0].text = str(int(row["Ano_Competencia"]))
            linha[1].text = str(int(row["Matricula"]))
            linha[2].text = f"R$ {float(row['Valor_Pago']):,.2f}"

        if not os.path.exists(PASTA_RELATORIOS):
            os.makedirs(PASTA_RELATORIOS)
        doc.add_paragraph("")
        doc.add_heading("Gráficos", level=2)
        doc.add_picture(gerar_grafico_valores(df_reemb), width=Inches(5.5))
        doc.add_picture(gerar_grafico_funcionarios(df_reemb), width=Inches(5.5))

        nome_arq = os.path.join(PASTA_RELATORIOS, f"Relatorio_Gerencial_{datetime.now():%Y%m%d_%H%M%S}.docx")
        doc.save(nome_arq)

        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, "[Bot] Relatório Word gerado com sucesso!\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)

    except Exception as e:
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, f"[Bot] Erro ao gerar relatório: {e}\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)

##Bloco 5

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm

def gerar_relatorio_pdf_gerencial():
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import cm
    import matplotlib.pyplot as plt

    try:
        # Conexão e leitura do banco
        conn = pyodbc.connect(f"DRIVER={{Microsoft Access Driver (*.mdb, *.accdb)}};DBQ={CAMINHO_BANCO};")
        df_reemb = pd.read_sql("SELECT * FROM ReembolsosEscolares", conn)
        df_func = pd.read_sql("SELECT * FROM Funcionarios", conn)
        conn.close()

        # Preparação dos dados
        df_reemb["Ano_Competencia"] = pd.to_numeric(df_reemb["Ano_Competencia"], errors="coerce").astype("Int64")
        df_reemb["Valor_Pago"] = df_reemb["Valor_Pago"].astype(str).str.replace(",", ".").astype(float)

        ativos_df = df_func[df_func["Status"].str.lower() == "ativo"]
        inativos_df = df_func[df_func["Status"].str.lower() != "ativo"]
        df_ativos = df_reemb[df_reemb["Matricula"].isin(ativos_df["Matricula"])]

        agrupado_todos = df_reemb.groupby("Ano_Competencia").agg(
            Funcionarios=("Matricula", "nunique"),
            Valor_Total=("Valor_Pago", "sum")
        ).reset_index()

        agrupado_ativos = df_ativos.groupby("Ano_Competencia").agg(
            Funcionarios_Ativos=("Matricula", "nunique"),
            Valor_Ativos=("Valor_Pago", "sum")
        ).reset_index()

        # Criação do PDF
        nome_arquivo = os.path.join(PASTA_RELATORIOS, f"Relatorio_Gerencial_{datetime.now():%Y%m%d_%H%M%S}.pdf")
        c = canvas.Canvas(nome_arquivo, pagesize=A4)
        largura, altura = A4
        y = altura - 2*cm

        # Título
        c.setFont("Helvetica-Bold", 16)
        c.drawString(2*cm, y, "Relatório Gerencial de Reembolsos")
        y -= 1*cm

        # Resumo de Funcionários
        c.setFont("Helvetica-Bold", 13)
        c.drawString(2*cm, y, "Resumo de Funcionários")
        y -= 0.6*cm
        c.setFont("Helvetica", 11)
        c.drawString(2*cm, y, f"Total: {len(df_func)}")
        y -= 0.5*cm
        c.drawString(2*cm, y, f"Ativos: {len(ativos_df)}")
        y -= 0.5*cm
        c.drawString(2*cm, y, f"Inativos: {len(inativos_df)}")
        y -= 1*cm

        # Totais por Ano (Todos)
        c.setFont("Helvetica-Bold", 13)
        c.drawString(2*cm, y, "Totais por Ano (Todos)")
        y -= 0.6*cm
        c.setFont("Helvetica-Bold", 11)
        c.drawString(2*cm, y, "Ano")
        c.drawString(6*cm, y, "Funcionários")
        c.drawString(11*cm, y, "Valor Total (R$)")
        y -= 0.4*cm
        c.setFont("Helvetica", 11)
        for _, row in agrupado_todos.iterrows():
            c.drawString(2*cm, y, str(row["Ano_Competencia"]))
            c.drawString(6*cm, y, str(row["Funcionarios"]))
            c.drawString(11*cm, y, f"R$ {row['Valor_Total']:,.2f}")
            y -= 0.5*cm
            if y < 5*cm:
                c.showPage()
                y = altura - 2*cm

        y -= 0.5*cm

        # Totais por Ano (Ativos)
        c.setFont("Helvetica-Bold", 13)
        c.drawString(2*cm, y, "Totais por Ano (Ativos)")
        y -= 0.6*cm
        c.setFont("Helvetica-Bold", 11)
        c.drawString(2*cm, y, "Ano")
        c.drawString(6*cm, y, "Funcionários Ativos")
        c.drawString(11*cm, y, "Valor Ativos (R$)")
        y -= 0.4*cm
        c.setFont("Helvetica", 11)
        for _, row in agrupado_ativos.iterrows():
            c.drawString(2*cm, y, str(row["Ano_Competencia"]))
            c.drawString(6*cm, y, str(row["Funcionarios_Ativos"]))
            c.drawString(11*cm, y, f"R$ {row['Valor_Ativos']:,.2f}")
            y -= 0.5*cm
            if y < 5*cm:
                c.showPage()
                y = altura - 2*cm

        # Gráficos
        fig1_path = os.path.join(PASTA_RELATORIOS, "grafico_valor_total_ano.png")
        df1 = df_reemb.groupby("Ano_Competencia")["Valor_Pago"].sum().reset_index()
        plt.figure(figsize=(6, 4))
        plt.bar(df1["Ano_Competencia"].astype(str), df1["Valor_Pago"], color="orange")
        plt.title("Valor Reembolsado por Ano")
        plt.xlabel("Ano")
        plt.ylabel("R$")
        plt.tight_layout()
        plt.savefig(fig1_path)
        plt.close()

        fig2_path = os.path.join(PASTA_RELATORIOS, "grafico_funcionarios_ano.png")
        df2 = df_reemb.groupby("Ano_Competencia")["Matricula"].nunique().reset_index()
        plt.figure(figsize=(6, 4))
        plt.bar(df2["Ano_Competencia"].astype(str), df2["Matricula"], color="green")
        plt.title("Funcionários Reembolsados por Ano")
        plt.xlabel("Ano")
        plt.ylabel("Qtd. Funcionários")
        plt.tight_layout()
        plt.savefig(fig2_path)
        plt.close()

        c.showPage()
        c.drawImage(fig1_path, 2*cm, altura/2, width=16*cm, preserveAspectRatio=True, mask='auto')
        c.drawImage(fig2_path, 2*cm, 2*cm, width=16*cm, preserveAspectRatio=True, mask='auto')
        c.showPage()
        c.save()

        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, "[Bot] Relatório PDF gerado com sucesso com todos os dados e gráficos!\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)

    except Exception as e:
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, f"[Bot] Erro ao gerar relatório PDF: {e}\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)

## Bloco 6

def gerar_apresentacao_powerpoint():
    try:
        from pptx import Presentation
        from pptx.util import Inches
        import matplotlib.pyplot as plt

        conn = pyodbc.connect(f"DRIVER={{Microsoft Access Driver (*.mdb, *.accdb)}};DBQ={CAMINHO_BANCO};")
        df_reemb = pd.read_sql("SELECT * FROM ReembolsosEscolares", conn)
        df_func = pd.read_sql("SELECT * FROM Funcionarios", conn)
        conn.close()

        df_reemb["Ano_Competencia"] = pd.to_numeric(df_reemb["Ano_Competencia"], errors="coerce").astype("Int64")
        df_reemb["Valor_Pago"] = df_reemb["Valor_Pago"].astype(str).str.replace(",", ".").astype(float)

        ativos_df = df_func[df_func["Status"].str.lower() == "ativo"]
        df_ativos = df_reemb[df_reemb["Matricula"].isin(ativos_df["Matricula"])]

        agrupado_todos = df_reemb.groupby("Ano_Competencia").agg(
            Funcionarios=("Matricula", "nunique"),
            Valor_Total=("Valor_Pago", "sum")
        ).reset_index()

        agrupado_ativos = df_ativos.groupby("Ano_Competencia").agg(
            Funcionarios_Ativos=("Matricula", "nunique"),
            Valor_Ativos=("Valor_Pago", "sum")
        ).reset_index()

        prs = Presentation()

        # Slide 1 - Título
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Relatório Gerencial de Reembolsos"
        subtitle.text = f"Gerado em: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}"

        # Slide 2 - Resumo
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Resumo de Funcionários"
        content = f"Total de funcionários: {len(df_func)}\nAtivos: {len(ativos_df)}\nInativos: {len(df_func) - len(ativos_df)}"
        slide.placeholders[1].text = content

        # Slide 3 - Totais por Ano (Todos)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Totais por Ano (Todos)"
        linhas = ["Ano | Funcionários | Valor Total (R$)"]
        for _, row in agrupado_todos.iterrows():
            linha = f"{int(row['Ano_Competencia'])} | {int(row['Funcionarios'])} | R$ {row['Valor_Total']:,.2f}"
            linhas.append(linha)
        slide.placeholders[1].text = "\n".join(linhas)

        # Slide 4 - Totais por Ano (Ativos)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Totais por Ano (Ativos)"
        linhas = ["Ano | Funcionários Ativos | Valor Ativos (R$)"]
        for _, row in agrupado_ativos.iterrows():
            linha = f"{int(row['Ano_Competencia'])} | {int(row['Funcionarios_Ativos'])} | R$ {row['Valor_Ativos']:,.2f}"
            linhas.append(linha)
        slide.placeholders[1].text = "\n".join(linhas)

        # Slide 5 - Gráfico de valor por ano
        fig1_path = os.path.join(PASTA_RELATORIOS, "grafico_valor_pptx.png")
        df1 = df_reemb.groupby("Ano_Competencia")["Valor_Pago"].sum().reset_index()
        plt.figure(figsize=(6, 4))
        plt.bar(df1["Ano_Competencia"].astype(str), df1["Valor_Pago"], color="orange")
        plt.title("Valor Reembolsado por Ano")
        plt.xlabel("Ano")
        plt.ylabel("Valor Total (R$)")
        plt.tight_layout()
        plt.savefig(fig1_path)
        plt.close()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Gráfico: Valor Reembolsado por Ano"
        slide.shapes.add_picture(fig1_path, Inches(1), Inches(1.5), width=Inches(8))

        # Slide 6 - Gráfico de funcionários por ano
        fig2_path = os.path.join(PASTA_RELATORIOS, "grafico_func_pptx.png")
        df2 = df_reemb.groupby("Ano_Competencia")["Matricula"].nunique().reset_index()
        plt.figure(figsize=(6, 4))
        plt.bar(df2["Ano_Competencia"].astype(str), df2["Matricula"], color="green")
        plt.title("Funcionários Reembolsados por Ano")
        plt.xlabel("Ano")
        plt.ylabel("Qtd. Funcionários")
        plt.tight_layout()
        plt.savefig(fig2_path)
        plt.close()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Gráfico: Funcionários Reembolsados por Ano"
        slide.shapes.add_picture(fig2_path, Inches(1), Inches(1.5), width=Inches(8))

        # Salvar
        nome_arq = os.path.join(PASTA_RELATORIOS, f"Apresentacao_Gerencial_{datetime.now():%Y%m%d_%H%M%S}.pptx")
        prs.save(nome_arq)

        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, "[Bot] Apresentação PowerPoint gerada com sucesso!\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)

    except Exception as e:
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, f"[Bot] Erro ao gerar apresentação: {e}\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)


# Bloco 4
def consultar_reembolsos():
    try:
        conn = pyodbc.connect(f"DRIVER={{Microsoft Access Driver (*.mdb, *.accdb)}};DBQ={CAMINHO_BANCO};")
        cursor = conn.cursor()
        cursor.execute("""
            SELECT Nome_Dependente, Mes_Competencia, Ano_Competencia, Valor_Pago, Tipo_Documento, Serie_Escolar
            FROM ReembolsosEscolares
            WHERE Matricula = ?
            ORDER BY Ano_Competencia DESC, Mes_Competencia DESC
        """, (matricula_usuario,))
        resultados = cursor.fetchall()
        conn.close()

        if not resultados:
            texto_chat.config(state="normal")
            texto_chat.insert(tk.END, "[Bot] Nenhum reembolso localizado para sua matrícula.\n")
            texto_chat.config(state="disabled")
            texto_chat.see(tk.END)
            return

        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, "[Bot] Reembolsos localizados:\n")
        for linha in resultados:
            nome_dep, mes, ano, valor, tipo, serie = linha
            try:
                mes = int(float(mes))
                ano = int(float(ano))
                valor = float(str(valor).replace(",", "."))
                resumo = f"- {nome_dep} | {mes:02d}/{ano} | R$ {valor:,.2f} | {tipo} | {serie}\n"
            except Exception as conv_erro:
                resumo = f"- Erro ao formatar linha: {conv_erro}\n"
            texto_chat.insert(tk.END, resumo)
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)

    except Exception as e:
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, f"[Bot] Erro ao consultar reembolsos: {e}\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)

def anexar_arquivo():
    global arquivos_reembolso
    caminho = filedialog.askopenfilename(title="Selecione o arquivo para upload")
    if not caminho: return
    nome = os.path.basename(caminho)
    tipo = classificar_por_conteudo(caminho)
    pastas = {
        "boleto": r"D:\visita\Documents\Sistema TCC\Chat Eletronuclear\Boleto",
        "comprovante": r"D:\visita\Documents\Sistema TCC\Chat Eletronuclear\Comprovantes de Pagamento Bancário",
        "formulario": r"D:\visita\Documents\Sistema TCC\Chat Eletronuclear\Formulário"
    }
    if tipo in pastas:
        shutil.copy2(caminho, os.path.join(pastas[tipo], nome))
        arquivos_reembolso[tipo] = caminho
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, f"[Bot] Arquivo '{nome}' classificado como '{tipo}' e salvo.\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)
    if len(arquivos_reembolso) == 3:
        dados = extrair_dados()
        if dados.get("valor_pago") and dados.get("nome_funcionario"):
            inserir_dados_no_access(dados)
            gerar_backup_excel()
        else:
            texto_chat.config(state="normal")
            texto_chat.insert(tk.END, "[Bot] Falha na extração de dados. Verifique se os arquivos estão legíveis.\n")
            texto_chat.config(state="disabled")
            texto_chat.see(tk.END)

def processar_mensagem(event=None):
    global fase
    texto = entrada_texto.get("1.0", tk.END).strip()
    entrada_texto.delete("1.0", tk.END)
    if not texto:
        return "break"
    texto_chat.config(state="normal")
    texto_chat.insert(tk.END, f"[Você] {texto}\n")
    texto_chat.config(state="disabled")
    texto_chat.see(tk.END)

    if texto.isdigit() and len(texto) >= 6:
        autenticar_matricula(texto)
    elif texto == "1":
        fase = "upload_reembolso"
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, "[Bot] Anexe boleto, comprovante e formulário em PDF.\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)
    elif texto == "2":
        consultar_reembolsos()
    elif texto == "3" and permissao_usuario.lower() == "administrador":
        gerar_relatorio_word_gerencial()
##Bloco 5
    elif texto == "4" and permissao_usuario.lower() == "administrador":
        gerar_relatorio_pdf_gerencial()

##Bloco 6

    elif texto == "5" and permissao_usuario.lower() == "administrador":
        gerar_apresentacao_powerpoint()


    else:
        texto_chat.config(state="normal")
        texto_chat.insert(tk.END, "[Bot] Opção inexistente ou não autorizada.\n")
        texto_chat.config(state="disabled")
        texto_chat.see(tk.END)
    return "break"

# Interface gráfica
janela = tk.Tk()
janela.title("Chat Eletronuclear")

if os.path.exists(LOGO_PATH):
    logo_img = Image.open(LOGO_PATH).resize((300, 60))
    logo_tk = ImageTk.PhotoImage(logo_img)
    tk.Label(janela, image=logo_tk).pack(pady=(5, 0))

frame_chat = tk.Frame(janela)
scrollbar = tk.Scrollbar(frame_chat)
texto_chat = tk.Text(frame_chat, height=20, width=60, yscrollcommand=scrollbar.set)
scrollbar.config(command=texto_chat.yview)
scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
texto_chat.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
frame_chat.pack(padx=10, pady=10)

texto_chat.insert(tk.END, "[Bot] Bem-vindo ao Chat Eletronuclear!\n")
texto_chat.insert(tk.END, "[Bot] Por favor, digite sua matrícula para continuar:\n")

entrada_texto = tk.Text(janela, height=1, width=60)
entrada_texto.pack(pady=(0, 10))
entrada_texto.bind("<Return>", processar_mensagem)
entrada_texto.focus_set()

botao_anexar = tk.Button(janela, text="📎 Anexar Arquivo", command=anexar_arquivo)
botao_anexar.pack(pady=5)

janela.mainloop()
